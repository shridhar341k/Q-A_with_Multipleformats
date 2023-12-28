import pandas as pd
import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from io import BytesIO
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
import os
import dotenv
openai_api_key = os.getenv("openai_api_key")

def get_pdf_text(pdf):
    text = ""
    pdf_reader = PdfReader(BytesIO(pdf.read()))
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def get_doc_text(doc):
    docx = Document(BytesIO(doc.read()))
    full_text = []
    for paragraph in docx.paragraphs:
        full_text.append(paragraph.text)
    return ' '.join(full_text)

def get_ppt_text(ppt):
    pres = Presentation(BytesIO(ppt.read()))
    full_text = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return ' '.join(full_text)



def get_data_text(data_files):
    text = ""
    for data_file in data_files:
        if data_file.type == "text/csv":
            df = pd.read_csv(data_file)
            text += ' '.join(df.astype(str).values.flatten())
        elif data_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            df = pd.read_excel(data_file)
            text += ' '.join(df.astype(str).values.flatten())
        elif data_file.type == "application/pdf":
            text += get_pdf_text(data_file)
        elif data_file.type == "application/msword" or data_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text += get_doc_text(data_file)
        elif data_file.type == "application/vnd.ms-powerpoint" or data_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text += get_ppt_text(data_file)
        else:
            st.error(f"File type not supported: {data_file.type}")
            return None
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI(openai_api_key=openai_api_key)
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)

def main():
    st.set_page_config(page_title="Chat with multiple data files",
                       page_icon=":books:")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.header("Chat with multiple data files :books:")
    user_question = st.text_input("Ask a question about your documents:")
    if user_question:
        handle_userinput(user_question)

    with st.sidebar:
        st.subheader("Your documents")
        data_files = st.file_uploader(
            "Upload your PDFs/CSVs/Excels/DOCs/PPTs/TIFs here and click on 'Process'",
            accept_multiple_files=True,
            type=['pdf', 'csv', 'xlsx', 'doc', 'docx', 'ppt', 'pptx']
        )
        if st.button("Process"):
            with st.spinner("Processing"):
                # get data text
                raw_text = get_data_text(data_files)

                # get the text chunks
                text_chunks = get_text_chunks(raw_text)

                # create vector store
                vectorstore = get_vectorstore(text_chunks)

                # create conversation chain
                st.session_state.conversation = get_conversation_chain(
                    vectorstore)

if __name__ == '__main__':
    main()
